from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import requests
from io import BytesIO
from content_generator import ContentGenerator

class PPTGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = Inches(16), Inches(9)
        self.content_generator = ContentGenerator()
        self.colors = {
            'primary': RGBColor(41, 128, 185),    # Blue
            'secondary': RGBColor(46, 204, 113),  # Green
            'accent': RGBColor(155, 89, 182),     # Purple
            'dark': RGBColor(44, 62, 80),         # Dark Blue
            'light': RGBColor(236, 240, 241)      # Light Gray
        }
        # Add decorative images
        self.decorative_images = {
            'tech_bg': 'https://images.unsplash.com/photo-1451187580459-43490279c0fa?w=800',
            'circuit': 'https://images.unsplash.com/photo-1518770660439-4636190af475?w=800',
            'innovation': 'https://images.unsplash.com/photo-1485827404703-89b55fcc595e?w=800'
        }
        # Add content-specific images
        self.content_images = {
            'news': 'https://images.unsplash.com/photo-1495020689067-958852a7765e?w=800',  # News background
            'tech': 'https://images.unsplash.com/photo-1518770660439-4636190af475?w=800',  # Tech circuit
            'quiz': 'https://images.unsplash.com/photo-1503676260728-1c00da094a0b?w=800',   # Quiz/education
            'joke': 'https://images.unsplash.com/photo-1485827404703-89b55fcc595e?w=800'    # Innovation/fun
        }

    def _add_slide_base(self, title_text=None):
        """Create base slide with clean background"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        # Add clean background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                  self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['light']
        
        if title_text:
            title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
            tf = title.text_frame
            tf.text = title_text
            tf.paragraphs[0].font.size = Pt(32)
            tf.paragraphs[0].font.color.rgb = self.colors['dark']
            tf.paragraphs[0].font.bold = True
        
        return slide

    def _add_text_box(self, slide, text, left, top, width, height, 
                     font_size=Pt(20), color='dark', bold=False, alignment=PP_ALIGN.LEFT):
        """Add formatted text box to slide"""
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.text = text
        p = tf.paragraphs[0]
        p.font.size = font_size
        p.font.color.rgb = self.colors[color]
        p.font.bold = bold
        p.alignment = alignment
        return box

    def _add_title_slide(self):
        """Add title slide with clean design"""
        slide = self._add_slide_base()
        current_date = datetime.now().strftime("%B %d, %Y")
        
        # Add main title with enhanced styling
        self._add_text_box(slide, "Daily Tech News Update", 2, 3, 12, 2,
                          font_size=Pt(56), color='dark', bold=True, alignment=PP_ALIGN.CENTER)
        
        # Add subtitle with decorative line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4), Inches(5),
            Inches(8), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors['primary']
        
        self._add_text_box(slide, f"Generated on {current_date}\nYour Daily Dose of Tech Innovation",
                          2, 5.5, 12, 1, font_size=Pt(32), color='dark', alignment=PP_ALIGN.CENTER)

    def _add_outline_slide(self):
        """Add outline slide showing presentation structure"""
        slide = self._add_slide_base()
        
        # Add main outline title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5),
            Inches(14), Inches(1.2)
        )
        tf = title_box.text_frame
        tf.text = "Today's Tech News Agenda"
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.colors['dark']
        p.alignment = PP_ALIGN.LEFT

        # Add decorative line under title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(1.8),
            Inches(14), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors['primary']
        
        # Add outline items with icons and frames
        outline_items = [
            ("📰 Latest Tech News Updates", "Stay updated with the most recent tech developments"),
            ("🎨 Tech Innovation of the Day", "Explore cutting-edge technological breakthroughs"),
            ("❓ Interactive Tech Quiz", "Test your knowledge with our tech challenge"),
            ("😄 Tech Humor Break", "End with a light-hearted tech joke")
        ]
        
        # Add each outline item with modern styling
        for i, (title, desc) in enumerate(outline_items):
            # Add item frame
            item_frame = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1), Inches(2.5 + i*1.5),
                Inches(14), Inches(1.2)
            )
            item_frame.fill.solid()
            item_frame.fill.fore_color.rgb = self.colors['light']
            item_frame.line.color.rgb = self.colors['primary']
            item_frame.line.width = Pt(1)
            
            # Add item title
            title_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(2.7 + i*1.5),
                Inches(13), Inches(0.4)
            )
            tf = title_box.text_frame
            tf.text = title
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = self.colors['dark']
            p.alignment = PP_ALIGN.LEFT
            
            # Add item description
            desc_box = slide.shapes.add_textbox(
                Inches(2), Inches(3.1 + i*1.5),
                Inches(12), Inches(0.4)
            )
            tf = desc_box.text_frame
            tf.text = desc
            p = tf.paragraphs[0]
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['primary']
            p.alignment = PP_ALIGN.LEFT

    def _add_ai_image_slide(self):
        """Add AI-generated image slide with stylish centered title"""
        slide = self._add_slide_base()
        
        # Add decorative line above title
        line_top = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6), Inches(0.8),
            Inches(4), Inches(0.05)
        )
        line_top.fill.solid()
        line_top.fill.fore_color.rgb = self.colors['primary']
        
        # Add main title with special formatting
        title_box = slide.shapes.add_textbox(
            Inches(2), Inches(1),
            Inches(12), Inches(1)
        )
        tf = title_box.text_frame
        tf.text = "TECH INNOVATION"
        p = tf.paragraphs[0]
        p.font.size = Pt(50)
        p.font.bold = True
        p.font.color.rgb = self.colors['dark']
        p.alignment = PP_ALIGN.CENTER
        
        # Add "OF THE DAY" subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(2), Inches(1.8),
            Inches(12), Inches(0.6)
        )
        tf = subtitle_box.text_frame
        tf.text = "OF THE DAY"
        p = tf.paragraphs[0]
        p.font.size = Pt(30)
        p.font.color.rgb = self.colors['primary']
        p.alignment = PP_ALIGN.CENTER
        
        # Add decorative line below title
        line_bottom = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6), Inches(2.5),
            Inches(4), Inches(0.05)
        )
        line_bottom.fill.solid()
        line_bottom.fill.fore_color.rgb = self.colors['primary']
        
        # Generate and add AI image
        image_url = self.content_generator.generate_tech_image()
        if image_url:
            try:
                response = requests.get(image_url)
                img_stream = BytesIO(response.content)
                
                # Add image with proper positioning
                slide.shapes.add_picture(
                    img_stream, 
                    Inches(3), Inches(3),
                    width=Inches(10),
                    height=Inches(5)
                )
            except Exception as e:
                print(f"Error adding AI image: {str(e)}")
                # Add descriptive text instead of image
                desc_box = slide.shapes.add_textbox(
                    Inches(3), Inches(3),
                    Inches(10), Inches(5)
                )
                tf = desc_box.text_frame
                tf.text = "Today's innovation focuses on emerging technologies in AI, machine learning, and sustainable tech solutions that are shaping our digital future."
                p = tf.paragraphs[0]
                p.font.size = Pt(28)
                p.font.color.rgb = self.colors['dark']
                p.alignment = PP_ALIGN.CENTER
                
        # Add decorative elements
        for i, color in enumerate(['primary', 'accent', 'secondary']):
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(1 + i*0.3), Inches(1 + i*0.3),
                Inches(0.5), Inches(0.5)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.colors[color]
            circle.line.color.rgb = self.colors['light']

    def _add_news_slide(self, article):
        """Add news slide with modern layout and image"""
        slide = self._add_slide_base()
        
        # Add title with modern styling
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5),
            Inches(14), Inches(1.2)
        )
        tf = title_box.text_frame
        tf.text = str(article.get('title', 'No Title Available'))
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.colors['dark']
        p.alignment = PP_ALIGN.LEFT

        # Add source tag in a modern pill shape
        source_frame = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), Inches(2),
            Inches(2.5), Inches(0.4)
        )
        source_frame.fill.solid()
        source_frame.fill.fore_color.rgb = self.colors['primary']
        source_frame.line.color.rgb = self.colors['primary']
        
        source_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(2.05),
            Inches(2), Inches(0.3)
        )
        tf = source_box.text_frame
        tf.text = str(article.get('source', 'Unknown'))
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = self.colors['light']
        p.alignment = PP_ALIGN.CENTER

        # Add "Read Full Article" button
        read_more_frame = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(12), Inches(2),
            Inches(3), Inches(0.4)
        )
        read_more_frame.fill.solid()
        read_more_frame.fill.fore_color.rgb = self.colors['secondary']
        read_more_frame.line.color.rgb = self.colors['secondary']
        
        read_more_box = slide.shapes.add_textbox(
            Inches(12.2), Inches(2.05),
            Inches(2.6), Inches(0.3)
        )
        tf = read_more_box.text_frame
        tf.text = "Read Full Article →"
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = self.colors['light']
        p.alignment = PP_ALIGN.CENTER

        # Try to add image and content in two columns, fallback to full-width content if image fails
        try:
            image_url = article.get('urlToImage', self.content_images['tech'])
            response = requests.get(image_url)
            img_stream = BytesIO(response.content)
            
            # If image loads successfully, use two-column layout
            left_content = slide.shapes.add_textbox(
                Inches(1), Inches(3),
                Inches(7), Inches(4)
            )
            tf = left_content.text_frame
            tf.text = str(article.get('description', 'No description available'))
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.color.rgb = self.colors['dark']
            p.alignment = PP_ALIGN.LEFT
            
            # Add image frame and image
            image_frame = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(9), Inches(3),
                Inches(6), Inches(4)
            )
            image_frame.fill.solid()
            image_frame.fill.fore_color.rgb = self.colors['light']
            image_frame.line.color.rgb = self.colors['primary']
            image_frame.line.width = Pt(2)
            
            image = slide.shapes.add_picture(
                img_stream,
                Inches(9.1), Inches(3.1),
                width=Inches(5.8),
                height=Inches(3.8)
            )
        except Exception as e:
            print(f"Error adding article image: {str(e)}")
            # If image fails, use full-width content layout
            content = slide.shapes.add_textbox(
                Inches(1), Inches(3),
                Inches(14), Inches(4)
            )
            tf = content.text_frame
            tf.text = str(article.get('description', 'No description available'))
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.color.rgb = self.colors['dark']
            p.alignment = PP_ALIGN.LEFT

        # Add bottom line separator
        bottom_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(7.5),
            Inches(14), Inches(0.03)
        )
        bottom_line.fill.solid()
        bottom_line.fill.fore_color.rgb = self.colors['primary']

    def _add_quiz_slides(self):
        """Add quiz slides with clean layout focused on content"""
        for i, quiz in enumerate(self.content_generator.generate_tech_quiz(2), 1):
            slide = self._add_slide_base()
            
            # Add title with modern styling
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(0.5),
                Inches(14), Inches(1.2)
            )
            tf = title_box.text_frame
            tf.text = f"Tech Quiz Challenge #{i}"
            p = tf.paragraphs[0]
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = self.colors['dark']
            p.alignment = PP_ALIGN.LEFT

            # Add quiz tag
            quiz_frame = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1), Inches(2),
                Inches(2), Inches(0.4)
            )
            quiz_frame.fill.solid()
            quiz_frame.fill.fore_color.rgb = self.colors['primary']
            quiz_frame.line.color.rgb = self.colors['primary']
            
            quiz_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(2.05),
                Inches(1.6), Inches(0.3)
            )
            tf = quiz_box.text_frame
            tf.text = "Question"
            p = tf.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = self.colors['light']
            p.alignment = PP_ALIGN.CENTER

            # Add main content area with question
            question_frame = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1), Inches(3),
                Inches(14), Inches(1.5)
            )
            question_frame.fill.solid()
            question_frame.fill.fore_color.rgb = self.colors['light']
            question_frame.line.color.rgb = self.colors['primary']
            question_frame.line.width = Pt(2)
            
            question_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(3.3),
                Inches(13), Inches(1)
            )
            tf = question_box.text_frame
            tf.text = quiz['question']
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = self.colors['dark']
            p.alignment = PP_ALIGN.LEFT

            # Add options in a grid layout
            for j, option in enumerate(quiz['options']):
                row = j // 2
                col = j % 2
                
                option_frame = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(1 + col*7.5), Inches(5 + row*1.2),
                    Inches(6.5), Inches(1)
                )
                option_frame.fill.solid()
                option_frame.fill.fore_color.rgb = self.colors['light']
                if option == quiz['correct']:
                    option_frame.line.color.rgb = self.colors['secondary']
                    option_frame.line.width = Pt(2)
                else:
                    option_frame.line.color.rgb = self.colors['primary']
                    option_frame.line.width = Pt(1)
                
                option_box = slide.shapes.add_textbox(
                    Inches(1.5 + col*7.5), Inches(5.2 + row*1.2),
                    Inches(5.5), Inches(0.6)
                )
                tf = option_box.text_frame
                tf.text = f"{chr(65+j)}. {option}"
                p = tf.paragraphs[0]
                p.font.size = Pt(20)
                if option == quiz['correct']:
                    p.font.color.rgb = self.colors['secondary']
                    p.font.bold = True
                else:
                    p.font.color.rgb = self.colors['dark']
                p.alignment = PP_ALIGN.LEFT

            # Add bottom line separator
            bottom_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1), Inches(7.5),
                Inches(14), Inches(0.03)
            )
            bottom_line.fill.solid()
            bottom_line.fill.fore_color.rgb = self.colors['primary']

    def _add_joke_slide(self):
        """Add joke slide with improved layout"""
        slide = self._add_slide_base()
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5),
            Inches(14), Inches(1)
        )
        tf = title_box.text_frame
        tf.text = "Tech Humor Break"
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        p.alignment = PP_ALIGN.LEFT
        
        # Add joke with proper spacing
        joke = self.content_generator.get_tech_joke()
        joke_box = slide.shapes.add_textbox(
            Inches(1), Inches(3),
            Inches(14), Inches(2)
        )
        tf = joke_box.text_frame
        tf.text = joke
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = self.colors['dark']
        p.alignment = PP_ALIGN.CENTER

    def generate_presentation(self, articles):
        """Generate complete presentation"""
        print("Generating presentation...")
        self._add_title_slide()
        self._add_outline_slide()  # Add outline slide after title
        self._add_ai_image_slide()
        self._add_quiz_slides()
        
        for article in articles[:15]:
            self._add_news_slide(article)
            
        self._add_joke_slide()
        
        # Save the presentation
        pptx_file = f"tech_news_{datetime.now().strftime('%Y%m%d')}.pptx"
        self.prs.save(pptx_file)
        return pptx_file 